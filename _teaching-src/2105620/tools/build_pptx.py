#!/usr/bin/env python3
"""Build a PowerPoint version of a Quarto/Reveal deck, with the spoken script
in the speaker-notes pane.

Each slide is the rendered Reveal slide captured at 2x as a full-bleed image, so
the PowerPoint looks exactly like the deck on the website — the Teal–Amber
theme, the bundled fonts, the layout, all of it. The speaker notes are the real
text from the `::: {.notes}` blocks in the .qmd.

The trade-off: text on the slides is a picture and cannot be edited in
PowerPoint. To change a slide, edit the .qmd and re-run this.

    python3 build_pptx.py w03

Requires: the private (non-public) render in 2105620-2026/_site, so the notes
are present in the HTML.
"""
import asyncio
import functools
import html as htmllib
import http.server
import pathlib
import re
import socketserver
import sys
import threading
from html.parser import HTMLParser

ROOT = pathlib.Path("/home/claude/2105620-2026")
SITE = ROOT / "_site"
OUT = pathlib.Path("/home/claude/files_out")
PORT = 8886
SCALE = 2  # 2560 x 1440 captures


# ---------------------------------------------------------------- notes
class NoteExtractor(HTMLParser):
    """Collect, in document order, the notes text of every top-level slide."""

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.slides = []          # one entry per <section>, "" when it has none
        self._depth = 0
        self._in_note = 0
        self._buf = []
        self._started = False

    def handle_starttag(self, tag, attrs):
        a = dict(attrs)
        if tag == "section":
            self._depth += 1
            if self._depth == 1:
                self.slides.append("")
                self._started = True
        elif self._started and tag == "aside" and "notes" in (a.get("class") or ""):
            self._in_note = 1
            self._buf = []
        elif self._in_note:
            if tag in ("p", "li", "tr", "div", "h1", "h2", "h3", "h4"):
                self._buf.append("\n")
            elif tag == "br":
                self._buf.append("\n")
            elif tag in ("td", "th") and self._buf and not self._buf[-1].endswith("\n"):
                self._buf.append("  ·  ")

    def handle_endtag(self, tag):
        if tag == "aside" and self._in_note:
            self._in_note = 0
            text = "".join(self._buf)
            text = re.sub(r"[ \t]+", " ", text)
            text = re.sub(r"\n[ \t]+", "\n", text)
            text = re.sub(r"\n{3,}", "\n\n", text).strip()
            if self.slides:
                self.slides[-1] = text
        elif tag == "section":
            self._depth = max(0, self._depth - 1)
        elif self._in_note and tag == "li":
            self._buf.append("\n")

    def handle_data(self, data):
        if self._in_note:
            self._buf.append(data)


def extract_notes(deck):
    src = (SITE / f"{deck}.html").read_text(encoding="utf-8")
    parser = NoteExtractor()
    parser.feed(src)
    return [htmllib.unescape(x) for x in parser.slides]


# ---------------------------------------------------------------- shots
async def capture(deck, shot_dir):
    from playwright.async_api import async_playwright

    shot_dir.mkdir(parents=True, exist_ok=True)
    async with async_playwright() as p:
        b = await p.chromium.launch()
        pg = await b.new_page(viewport={"width": 1280, "height": 720},
                              device_scale_factor=SCALE)
        await pg.goto(f"http://127.0.0.1:{PORT}/{deck}.html", wait_until="networkidle")
        await pg.wait_for_timeout(2000)
        # the on-screen navigation would otherwise be baked into every picture
        await pg.add_style_tag(content=(
            ".reveal .controls, .reveal .progress, .reveal .slide-number,"
            " .reveal .speaker-controls { display: none !important; }"))
        await pg.wait_for_timeout(300)
        n = await pg.evaluate("Reveal.getTotalSlides()")
        for i in range(n):
            await pg.evaluate(f"Reveal.slide({i})")
            await pg.wait_for_timeout(280)
            await pg.screenshot(path=str(shot_dir / f"{i + 1:02d}.png"))
        await b.close()
    return n


def serve():
    handler = functools.partial(http.server.SimpleHTTPRequestHandler,
                                directory=str(SITE))

    class Q(socketserver.TCPServer):
        allow_reuse_address = True

    srv = Q(("127.0.0.1", PORT), handler)
    threading.Thread(target=srv.serve_forever, daemon=True).start()
    return srv


# ---------------------------------------------------------------- build
def build(deck, shot_dir, notes, title):
    from pptx import Presentation
    from pptx.util import Inches, Emu

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]           # completely empty layout

    shots = sorted(shot_dir.glob("*.png"))
    for i, img in enumerate(shots):
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(img), Emu(0), Emu(0),
                             width=prs.slide_width, height=prs.slide_height)
        note = notes[i] if i < len(notes) else ""
        if note:
            s.notes_slide.notes_text_frame.text = note

    prs.core_properties.title = title
    out = OUT / f"{deck}_slides_with_script.pptx"
    prs.save(str(out))
    return out, len(shots)


def main():
    deck = sys.argv[1]
    titles = {
        "w01": "2105620 Week 1 — Course framing and the publication landscape",
        "w02": "2105620 Week 2 — Research integrity, authorship, and generative AI",
        "w03": "2105620 Week 3 — Efficient literature review and evidence mapping",
    }
    notes = extract_notes(deck)
    srv = serve()
    shot_dir = pathlib.Path(f"/home/claude/pptx_shots/{deck}")
    try:
        n = asyncio.run(capture(deck, shot_dir))
    finally:
        srv.shutdown()

    out, made = build(deck, shot_dir, notes, titles.get(deck, deck))
    with_notes = sum(1 for x in notes if x)
    print(f"{deck}: {made} slides · {with_notes} of {len(notes)} carry a script "
          f"· {out.stat().st_size // 1024} KB")
    if made != n:
        print(f"  WARNING: captured {made} images for {n} reported slides")
    if len(notes) != made:
        print(f"  WARNING: {len(notes)} note slots for {made} slides")


if __name__ == "__main__":
    main()
