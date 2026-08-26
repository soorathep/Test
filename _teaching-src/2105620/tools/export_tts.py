#!/usr/bin/env python3
"""Export the spoken script of a deck as one plain-text file per slide, ready to
paste or batch into a text-to-speech engine.

The text is the ::: {.notes} narration with every trace of markup removed. No
stage directions, no timing cues, no brackets — a TTS engine reads what it is
given, so anything that is not meant to be heard must not be in the file.

    python3 export_tts.py w04
"""
import csv, pathlib, re, sys, unicodedata
from pptx import Presentation

DECK = sys.argv[1] if len(sys.argv) > 1 else "w04"
OUT = pathlib.Path(f"/home/claude/files_out/{DECK}_tts")
WPM = 140  # measured speaking rate for this narration style

# which slides belong to which recording
SEGMENTS = [
    (1, 10, "segment1_what-each-section-is-for"),
    (11, 19, "segment2_the-reverse-outline"),
    (20, 30, "segment3_where-arguments-break"),
]

# characters a TTS engine either mispronounces or reads aloud as punctuation
REPL = {
    "—": ", ", "–": "-", "‘": "'", "’": "'",
    "“": "", "”": "", "·": ",", "…": ".",
    " ": " ", "&": " and ",
}


def clean(t):
    t = unicodedata.normalize("NFC", t)
    for a, b in REPL.items():
        t = t.replace(a, b)
    t = re.sub(r"\*+", "", t)               # stray markdown emphasis
    t = re.sub(r"\[[^\]]*\]", "", t)        # any bracketed cue
    t = re.sub(r"[ \t]+", " ", t)
    t = re.sub(r" *\n *", "\n", t)
    t = re.sub(r"\n{3,}", "\n\n", t)
    t = re.sub(r" +([,.;:?!])", r"\1", t)
    t = re.sub(r",\s*,", ",", t)
    return t.strip()


def title_of(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().split("\n")[0]
    return ""


def main():
    prs = Presentation(f"/home/claude/files_out/{DECK}_slides_with_script.pptx")
    slides = list(prs.slides)
    if OUT.exists():
        for f in OUT.rglob("*"):
            if f.is_file():
                f.unlink()
    rows, total = [], 0
    for lo, hi, folder in SEGMENTS:
        d = OUT / folder
        d.mkdir(parents=True, exist_ok=True)
        seg_words = 0
        for n in range(lo, hi + 1):
            s = slides[n - 1]
            raw = s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
            text = clean(raw)
            words = len(text.split())
            seg_words += words
            f = d / f"s{n:02d}.txt"
            f.write_text(text + ("\n" if text else ""), encoding="utf-8")
            rows.append([n, folder.split("_")[0], words, round(words / WPM * 60),
                         f"{folder}/{f.name}"])
        total += seg_words
        print(f"{folder}: slides {lo}-{hi} · {seg_words} words "
              f"· about {seg_words / WPM:.1f} min")

    with (OUT / "manifest.csv").open("w", newline="", encoding="utf-8") as fh:
        w = csv.writer(fh)
        w.writerow(["slide", "segment", "words", "seconds_at_140wpm", "file"])
        w.writerows(rows)

    readme = OUT / "README.txt"
    readme.write_text(
        "Spoken script for 2105620 Week 4, one plain-text file per slide.\n\n"
        "Folders are the three recordings. File s07.txt is the narration for\n"
        "slide 7 of w04_slides.pdf and of w04_slides_with_script.pptx — the\n"
        "numbering is the same in all three.\n\n"
        "The text is already stripped of markup, stage directions and timing\n"
        "cues, so it can go straight into a text-to-speech engine. Estimated\n"
        f"durations in manifest.csv assume {WPM} words per minute.\n\n"
        "Slide 1 is the title card and has no narration; its file is empty.\n",
        encoding="utf-8")
    print(f"total {total} words · about {total / WPM:.0f} min of speech")
    print(f"-> {OUT}")


if __name__ == "__main__":
    main()
